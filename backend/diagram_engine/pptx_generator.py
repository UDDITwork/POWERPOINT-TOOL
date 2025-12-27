"""
PowerPoint Diagram Generator using python-pptx

This module generates fully editable PPTX files from structured JSON specifications.
Every shape, arrow, and text box is created as a native PowerPoint object,
ensuring complete editability when opened in PowerPoint.

Author: AI Patent Diagram Generator
License: MIT
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE, MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from typing import Dict, List, Any, Optional, Tuple
import logging

logger = logging.getLogger(__name__)


class PPTXDiagramGenerator:
    """
    Main class for generating PowerPoint diagrams from JSON specifications.

    Supports:
    - 182 AutoShapes
    - Connectors and arrows
    - Tables
    - Charts
    - Text formatting
    - Grouping and positioning
    """

    # Map of common shape names to MSO_SHAPE constants
    SHAPE_TYPE_MAP = {
        # Basic Shapes
        'rectangle': MSO_SHAPE.RECTANGLE,
        'rounded_rectangle': MSO_SHAPE.ROUNDED_RECTANGLE,
        'roundedrectangle': MSO_SHAPE.ROUNDED_RECTANGLE,  # Frontend camelCase
        'oval': MSO_SHAPE.OVAL,
        'circle': MSO_SHAPE.OVAL,
        'square': MSO_SHAPE.RECTANGLE,  # Square is just a rectangle
        'diamond': MSO_SHAPE.DIAMOND,
        'triangle': MSO_SHAPE.ISOSCELES_TRIANGLE,
        'triangledown': MSO_SHAPE.ISOSCELES_TRIANGLE,  # Will need rotation
        'inverted_triangle': MSO_SHAPE.ISOSCELES_TRIANGLE,
        'right_triangle': MSO_SHAPE.RIGHT_TRIANGLE,
        'parallelogram': MSO_SHAPE.PARALLELOGRAM,
        'trapezoid': MSO_SHAPE.TRAPEZOID,
        'hexagon': MSO_SHAPE.HEXAGON,
        'octagon': MSO_SHAPE.OCTAGON,
        'pentagon': MSO_SHAPE.PENTAGON,
        'cross': MSO_SHAPE.CROSS,
        'star': MSO_SHAPE.STAR_5_POINT,
        'heart': MSO_SHAPE.HEART,
        'cube': MSO_SHAPE.CUBE,

        # Flowchart Shapes (CRITICAL FOR PATENTS)
        'process': MSO_SHAPE.FLOWCHART_PROCESS,
        'flowchart_process': MSO_SHAPE.FLOWCHART_PROCESS,
        'decision': MSO_SHAPE.FLOWCHART_DECISION,
        'flowchart_decision': MSO_SHAPE.FLOWCHART_DECISION,
        'terminator': MSO_SHAPE.FLOWCHART_TERMINATOR,
        'flowchart_terminator': MSO_SHAPE.FLOWCHART_TERMINATOR,
        'data': MSO_SHAPE.FLOWCHART_DATA,
        'flowchart_data': MSO_SHAPE.FLOWCHART_DATA,
        'document': MSO_SHAPE.FLOWCHART_DOCUMENT,
        'flowchart_document': MSO_SHAPE.FLOWCHART_DOCUMENT,
        'multidocument': MSO_SHAPE.FLOWCHART_MULTIDOCUMENT,
        'predefined_process': MSO_SHAPE.FLOWCHART_PREDEFINED_PROCESS,
        'predefinedprocess': MSO_SHAPE.FLOWCHART_PREDEFINED_PROCESS,
        'flowchart_preparation': MSO_SHAPE.FLOWCHART_PREPARATION,
        'preparation': MSO_SHAPE.FLOWCHART_PREPARATION,
        'flowchart_manual_input': MSO_SHAPE.FLOWCHART_MANUAL_INPUT,
        'manualinput': MSO_SHAPE.FLOWCHART_MANUAL_INPUT,
        'flowchart_manual_operation': MSO_SHAPE.FLOWCHART_MANUAL_OPERATION,
        'manualoperation': MSO_SHAPE.FLOWCHART_MANUAL_OPERATION,
        'flowchart_connector': MSO_SHAPE.FLOWCHART_CONNECTOR,
        'flowchart_delay': MSO_SHAPE.FLOWCHART_DELAY,
        'delay': MSO_SHAPE.FLOWCHART_DELAY,
        'flowchart_merge': MSO_SHAPE.FLOWCHART_MERGE,
        'merge': MSO_SHAPE.FLOWCHART_MERGE,
        'extract': MSO_SHAPE.FLOWCHART_EXTRACT,
        'flowchart_or': MSO_SHAPE.FLOWCHART_OR,
        'or': MSO_SHAPE.FLOWCHART_OR,
        'summing': MSO_SHAPE.FLOWCHART_SUMMING_JUNCTION,
        'flowchart_stored_data': MSO_SHAPE.FLOWCHART_STORED_DATA,
        'storeddata': MSO_SHAPE.FLOWCHART_STORED_DATA,
        'internalstorage': MSO_SHAPE.FLOWCHART_INTERNAL_STORAGE,
        'offpageconnector': MSO_SHAPE.FLOWCHART_OFFPAGE_CONNECTOR,

        # Arrows
        'left_arrow': MSO_SHAPE.LEFT_ARROW,
        'arrowleft': MSO_SHAPE.LEFT_ARROW,
        'right_arrow': MSO_SHAPE.RIGHT_ARROW,
        'arrowright': MSO_SHAPE.RIGHT_ARROW,
        'up_arrow': MSO_SHAPE.UP_ARROW,
        'arrowup': MSO_SHAPE.UP_ARROW,
        'down_arrow': MSO_SHAPE.DOWN_ARROW,
        'arrowdown': MSO_SHAPE.DOWN_ARROW,
        'left_right_arrow': MSO_SHAPE.LEFT_RIGHT_ARROW,
        'arrowbidirectional': MSO_SHAPE.LEFT_RIGHT_ARROW,
        'up_down_arrow': MSO_SHAPE.UP_DOWN_ARROW,
        'arrowvertical': MSO_SHAPE.UP_DOWN_ARROW,
        'quad_arrow': MSO_SHAPE.QUAD_ARROW,
        'bent_arrow': MSO_SHAPE.BENT_ARROW,
        'arrowbent': MSO_SHAPE.BENT_ARROW,
        'curved_right_arrow': MSO_SHAPE.CURVED_RIGHT_ARROW,
        'arrowcurved': MSO_SHAPE.CURVED_RIGHT_ARROW,
        'curved_left_arrow': MSO_SHAPE.CURVED_LEFT_ARROW,
        'circular_arrow': MSO_SHAPE.CIRCULAR_ARROW,
        'arrowcircular': MSO_SHAPE.CIRCULAR_ARROW,
        'arrowuturn': MSO_SHAPE.U_TURN_ARROW,
        'chevron': MSO_SHAPE.CHEVRON,
        'chevronright': MSO_SHAPE.CHEVRON,
        'chevronleft': MSO_SHAPE.CHEVRON,
        'arrowpentagon': MSO_SHAPE.PENTAGON,
        'arrownotched': MSO_SHAPE.NOTCHED_RIGHT_ARROW,
        'arrowstriped': MSO_SHAPE.STRIPED_RIGHT_ARROW,

        # Callouts (for annotations)
        'callout': MSO_SHAPE.RECTANGULAR_CALLOUT,
        'rectangular_callout': MSO_SHAPE.RECTANGULAR_CALLOUT,
        'calloutrect': MSO_SHAPE.RECTANGULAR_CALLOUT,
        'rounded_rectangular_callout': MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT,
        'calloutrounded': MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT,
        'rounded_callout': MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT,
        'oval_callout': MSO_SHAPE.OVAL_CALLOUT,
        'calloutoval': MSO_SHAPE.OVAL_CALLOUT,
        'cloud_callout': MSO_SHAPE.CLOUD_CALLOUT,
        'calloutcloud': MSO_SHAPE.CLOUD_CALLOUT,
        'line_callout': MSO_SHAPE.LINE_CALLOUT_1,
        'calloutline': MSO_SHAPE.LINE_CALLOUT_1,
        'annotation': MSO_SHAPE.LINE_CALLOUT_1,
        'bracket': MSO_SHAPE.LEFT_BRACKET,
        'brace': MSO_SHAPE.LEFT_BRACE,

        # Containers
        'cylinder': MSO_SHAPE.CAN,
        'database': MSO_SHAPE.CAN,
        'can': MSO_SHAPE.CAN,
        'folder': MSO_SHAPE.FOLDED_CORNER,  # No FOLDER in python-pptx, using folded corner
        'card': MSO_SHAPE.FOLDED_CORNER,
        'frame': MSO_SHAPE.FRAME,
        'tape': MSO_SHAPE.FLOWCHART_PUNCHED_TAPE,
        'cube': MSO_SHAPE.CUBE,

        # Stars
        'star_4': MSO_SHAPE.STAR_4_POINT,
        'star4': MSO_SHAPE.STAR_4_POINT,
        'star_4_point': MSO_SHAPE.STAR_4_POINT,
        'star_5': MSO_SHAPE.STAR_5_POINT,
        'star_6': MSO_SHAPE.STAR_6_POINT,
        'star_8': MSO_SHAPE.STAR_8_POINT,
        'star_16': MSO_SHAPE.STAR_16_POINT,
        'star_24': MSO_SHAPE.STAR_24_POINT,

        # Math Symbols
        'plus': MSO_SHAPE.MATH_PLUS,
        'minus': MSO_SHAPE.MATH_MINUS,
        'multiply': MSO_SHAPE.MATH_MULTIPLY,
        'divide': MSO_SHAPE.MATH_DIVIDE,
        'equal': MSO_SHAPE.MATH_EQUAL,

        # Special / Symbols
        'cloud': MSO_SHAPE.CLOUD,
        'lightning': MSO_SHAPE.LIGHTNING_BOLT,
        'sun': MSO_SHAPE.SUN,
        'moon': MSO_SHAPE.MOON,
        'gear': MSO_SHAPE.GEAR_6,
        'smiley': MSO_SHAPE.SMILEY_FACE,
        'no_symbol': MSO_SHAPE.NO_SYMBOL,
        'shield': MSO_SHAPE.FLOWCHART_DISPLAY,  # Closest match
        'checkmark': MSO_SHAPE.RECTANGLE,  # Fallback
        'xmark': MSO_SHAPE.RECTANGLE,  # Fallback
        'warning': MSO_SHAPE.ISOSCELES_TRIANGLE,  # Warning triangle
        'info': MSO_SHAPE.OVAL,  # Info circle
        'flag': MSO_SHAPE.WAVE,  # Closest match
        'puzzle': MSO_SHAPE.RECTANGLE,  # Fallback
        'lock': MSO_SHAPE.RECTANGLE,  # Fallback
        'key': MSO_SHAPE.RECTANGLE,  # Fallback

        # Lines (rendered as thin rectangles or connectors)
        'line': MSO_SHAPE.RECTANGLE,
        'linearrow': MSO_SHAPE.RIGHT_ARROW,
        'linedouble': MSO_SHAPE.LEFT_RIGHT_ARROW,
    }

    # Connector types
    CONNECTOR_TYPE_MAP = {
        'straight': MSO_CONNECTOR.STRAIGHT,
        'elbow': MSO_CONNECTOR.ELBOW,
        'curve': MSO_CONNECTOR.CURVE,
        'curved': MSO_CONNECTOR.CURVE,
    }

    def __init__(self, width_inches: float = 10.0, height_inches: float = 7.5):
        """
        Initialize the PPTX generator.

        Args:
            width_inches: Slide width in inches (default: 10.0)
            height_inches: Slide height in inches (default: 7.5)
        """
        self.prs = Presentation()
        self.prs.slide_width = Inches(width_inches)
        self.prs.slide_height = Inches(height_inches)
        self.shape_registry: Dict[str, Any] = {}

        logger.info(f"Initialized PPTX generator with dimensions: {width_inches}x{height_inches} inches")

    def create_from_json(self, spec: Dict[str, Any]) -> Presentation:
        """
        Generate a complete presentation from JSON specification.

        Args:
            spec: Dictionary containing diagram specification
                {
                    "metadata": {"title": "...", "author": "..."},
                    "slides": [
                        {
                            "elements": [...],
                            "layout": {...}
                        }
                    ]
                }

        Returns:
            Presentation object ready to be saved
        """
        logger.info(f"Creating presentation from JSON spec")

        # Extract metadata
        metadata = spec.get('metadata', {})
        if 'title' in metadata:
            self.prs.core_properties.title = metadata['title']
        if 'author' in metadata:
            self.prs.core_properties.author = metadata['author']

        # Process slides
        slides_data = spec.get('slides', [spec])  # Support single slide or multiple

        for slide_idx, slide_data in enumerate(slides_data):
            logger.info(f"Processing slide {slide_idx + 1}/{len(slides_data)}")
            self._create_slide(slide_data)

        logger.info(f"Presentation created successfully with {len(self.prs.slides)} slides")
        return self.prs

    def _create_slide(self, slide_data: Dict[str, Any]) -> None:
        """Create a single slide from specification."""
        # Use blank layout
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Clear shape registry for this slide
        self.shape_registry.clear()

        elements = slide_data.get('elements', [])

        # Two-pass rendering:
        # Pass 1: Create all shapes and register them
        # Pass 2: Create connectors (which reference shapes)

        shapes_to_create = [e for e in elements if e.get('type') != 'connector']
        connectors_to_create = [e for e in elements if e.get('type') == 'connector']

        logger.debug(f"Creating {len(shapes_to_create)} shapes and {len(connectors_to_create)} connectors")

        # Pass 1: Shapes
        for element in shapes_to_create:
            self._add_element(slide, element)

        # Pass 2: Connectors
        for element in connectors_to_create:
            self._add_connector(slide, element)

    def _add_element(self, slide, element: Dict[str, Any]) -> Any:
        """
        Add an element to the slide based on its type.

        Supported types:
        - shape: Any of the 182 AutoShapes
        - textbox: Standalone text box
        - table: Data table
        - chart: Data chart
        - image: Picture
        """
        element_type = element.get('type', 'rectangle')
        element_id = element.get('id', f"elem_{len(self.shape_registry)}")

        logger.debug(f"Adding element: {element_id} (type: {element_type})")

        if element_type == 'textbox':
            shape = self._add_textbox(slide, element)
        elif element_type == 'table':
            shape = self._add_table(slide, element)
        elif element_type == 'chart':
            shape = self._add_chart(slide, element)
        elif element_type == 'image':
            shape = self._add_image(slide, element)
        else:
            # Default: treat as shape
            shape = self._add_shape(slide, element)

        # Register shape for connector references
        if element_id:
            self.shape_registry[element_id] = shape

        return shape

    def _add_shape(self, slide, element: Dict[str, Any]) -> Any:
        """Add a shape to the slide."""
        shape_type = element.get('type', 'rectangle')

        # Get MSO_SHAPE constant - normalize shape name
        shape_type_normalized = shape_type.lower().replace('_', '').replace('-', '')
        mso_shape = self.SHAPE_TYPE_MAP.get(shape_type_normalized)

        if mso_shape is None:
            # Try original lowercase
            mso_shape = self.SHAPE_TYPE_MAP.get(shape_type.lower())

        if mso_shape is None:
            logger.warning(f"Unknown shape type '{shape_type}', defaulting to rectangle")
            mso_shape = MSO_SHAPE.RECTANGLE

        # Position and size with validation
        pos = element.get('position', {'x': 1.0, 'y': 1.0})
        size = element.get('size', {'width': 2.0, 'height': 1.0})

        # Ensure position values are valid numbers
        try:
            x = float(pos.get('x', 1.0))
            y = float(pos.get('y', 1.0))
            width = float(size.get('width', 2.0))
            height = float(size.get('height', 1.0))
        except (TypeError, ValueError) as e:
            logger.warning(f"Invalid position/size values for element, using defaults: {e}")
            x, y, width, height = 1.0, 1.0, 2.0, 1.0

        # Clamp to reasonable bounds
        x = max(0.1, min(9.5, x))
        y = max(0.1, min(7.0, y))
        width = max(0.3, min(9.0, width))
        height = max(0.3, min(6.5, height))

        try:
            shape = slide.shapes.add_shape(
                mso_shape,
                Inches(x),
                Inches(y),
                Inches(width),
                Inches(height)
            )
        except Exception as e:
            logger.error(f"Failed to add shape '{shape_type}': {e}")
            # Fallback to rectangle
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x),
                Inches(y),
                Inches(width),
                Inches(height)
            )

        # Apply styling
        self._apply_shape_styling(shape, element)

        # Add text if provided
        if 'text' in element:
            self._add_text_to_shape(shape, element)

        return shape

    def _add_textbox(self, slide, element: Dict[str, Any]) -> Any:
        """Add a standalone text box."""
        pos = element.get('position', {'x': 1.0, 'y': 1.0})
        size = element.get('size', {'width': 2.0, 'height': 1.0})

        textbox = slide.shapes.add_textbox(
            Inches(pos['x']),
            Inches(pos['y']),
            Inches(size['width']),
            Inches(size['height'])
        )

        text_frame = textbox.text_frame
        text_frame.text = element.get('text', '')

        # Text formatting
        self._apply_text_formatting(text_frame, element)

        return textbox

    def _add_connector(self, slide, element: Dict[str, Any]) -> Any:
        """
        Add a connector between two shapes.

        Element format:
        {
            "type": "connector",
            "connector_type": "straight|elbow|curve",
            "from": "shape_id",
            "to": "shape_id",
            "from_side": "top|bottom|left|right|center",
            "to_side": "top|bottom|left|right|center",
            "label": "optional text",
            "style": {...}
        }
        """
        from_id = element.get('from')
        to_id = element.get('to')
        connector_id = element.get('id', 'unknown')

        if not from_id or not to_id:
            logger.error(f"Connector '{connector_id}' missing from/to references: {element}")
            # Don't fail silently - log error but continue to render other elements
            return None

        from_shape = self.shape_registry.get(from_id)
        to_shape = self.shape_registry.get(to_id)

        missing_shapes = []
        if not from_shape:
            missing_shapes.append(f"source '{from_id}'")
        if not to_shape:
            missing_shapes.append(f"target '{to_id}'")

        if missing_shapes:
            # Log detailed error for debugging
            logger.error(f"Connector '{connector_id}' references non-existent shape(s): {', '.join(missing_shapes)}")
            logger.error(f"Available shapes in registry: {list(self.shape_registry.keys())}")
            # Return None to skip this connector but continue rendering
            return None

        # Auto-detect optimal connection sides based on shape positions
        from_side = element.get('from_side')
        to_side = element.get('to_side')

        if not from_side or not to_side:
            from_side, to_side = self._calculate_optimal_connection_sides(from_shape, to_shape)

        # Calculate connection points
        from_point = self._get_connection_point(from_shape, from_side)
        to_point = self._get_connection_point(to_shape, to_side)

        # Connector type
        connector_type_str = element.get('connector_type', 'straight')
        connector_type = self.CONNECTOR_TYPE_MAP.get(connector_type_str, MSO_CONNECTOR.STRAIGHT)

        # Create connector
        connector = slide.shapes.add_connector(
            connector_type,
            from_point[0], from_point[1],
            to_point[0], to_point[1]
        )

        # Apply connector styling
        self._apply_connector_styling(connector, element)

        return connector

    def _calculate_optimal_connection_sides(self, from_shape, to_shape) -> Tuple[str, str]:
        """
        Calculate optimal connection sides based on relative positions of shapes.

        Returns:
            Tuple of (from_side, to_side)
        """
        # Get shape centers
        from_cx = from_shape.left + from_shape.width // 2
        from_cy = from_shape.top + from_shape.height // 2
        to_cx = to_shape.left + to_shape.width // 2
        to_cy = to_shape.top + to_shape.height // 2

        # Calculate relative position
        dx = to_cx - from_cx
        dy = to_cy - from_cy

        # Determine primary direction
        if abs(dy) > abs(dx):
            # Vertical arrangement
            if dy > 0:
                # Target is below source
                return 'bottom', 'top'
            else:
                # Target is above source
                return 'top', 'bottom'
        else:
            # Horizontal arrangement
            if dx > 0:
                # Target is to the right
                return 'right', 'left'
            else:
                # Target is to the left
                return 'left', 'right'

    def _add_table(self, slide, element: Dict[str, Any]) -> Any:
        """Add a table to the slide."""
        pos = element.get('position', {'x': 1.0, 'y': 1.0})
        size = element.get('size', {'width': 6.0, 'height': 3.0})

        rows = element.get('rows', 3)
        cols = element.get('cols', 3)

        table_shape = slide.shapes.add_table(
            rows, cols,
            Inches(pos['x']),
            Inches(pos['y']),
            Inches(size['width']),
            Inches(size['height'])
        )

        table = table_shape.table

        # Populate table data if provided
        data = element.get('data', [])
        for row_idx, row_data in enumerate(data):
            if row_idx >= rows:
                break
            for col_idx, cell_value in enumerate(row_data):
                if col_idx >= cols:
                    break
                table.cell(row_idx, col_idx).text = str(cell_value)

        return table_shape

    def _add_chart(self, slide, element: Dict[str, Any]) -> Any:
        """Add a chart to the slide."""
        pos = element.get('position', {'x': 1.0, 'y': 1.0})
        size = element.get('size', {'width': 6.0, 'height': 4.0})

        chart_type_str = element.get('chart_type', 'column')
        chart_type_map = {
            'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
            'bar': XL_CHART_TYPE.BAR_CLUSTERED,
            'line': XL_CHART_TYPE.LINE,
            'pie': XL_CHART_TYPE.PIE,
            'scatter': XL_CHART_TYPE.XY_SCATTER,
        }
        chart_type = chart_type_map.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

        # Chart data
        chart_data = CategoryChartData()
        categories = element.get('categories', ['A', 'B', 'C'])
        chart_data.categories = categories

        series_list = element.get('series', [{'name': 'Series 1', 'values': [1, 2, 3]}])
        for series in series_list:
            chart_data.add_series(series['name'], series['values'])

        chart_shape = slide.shapes.add_chart(
            chart_type,
            Inches(pos['x']),
            Inches(pos['y']),
            Inches(size['width']),
            Inches(size['height']),
            chart_data
        )

        return chart_shape

    def _add_image(self, slide, element: Dict[str, Any]) -> Any:
        """Add an image to the slide."""
        image_path = element.get('path')
        if not image_path:
            logger.warning("Image element missing path")
            return None

        pos = element.get('position', {'x': 1.0, 'y': 1.0})
        size = element.get('size', {'width': 3.0, 'height': 2.0})

        picture = slide.shapes.add_picture(
            image_path,
            Inches(pos['x']),
            Inches(pos['y']),
            width=Inches(size['width']),
            height=Inches(size['height'])
        )

        return picture

    def _apply_shape_styling(self, shape, element: Dict[str, Any]) -> None:
        """Apply fill, line, and other styling to a shape.

        For patent diagrams, default is:
        - Transparent fill (no fill)
        - Black border (1pt)
        - Solid line style
        """
        style = element.get('style', {})

        # Fill color - default to NO FILL (transparent) for patent diagrams
        if 'fill_color' in style and style['fill_color']:
            fill_color = style['fill_color']
            shape.fill.solid()
            if isinstance(fill_color, str):
                rgb = self._hex_to_rgb(fill_color)
                shape.fill.fore_color.rgb = RGBColor(*rgb)
        else:
            # No fill = transparent background
            shape.fill.background()

        # Line color - default to BLACK
        line_color = style.get('line_color', '000000')
        rgb = self._hex_to_rgb(line_color)
        shape.line.color.rgb = RGBColor(*rgb)

        # Line width - default to 1pt
        line_width = style.get('line_width', 1.0)
        shape.line.width = Pt(line_width)

        # Line style (dashed, dotted, etc.) - default to solid
        if 'line_style' in style:
            line_style_map = {
                'solid': MSO_LINE_DASH_STYLE.SOLID,
                'dash': MSO_LINE_DASH_STYLE.DASH,
                'dot': MSO_LINE_DASH_STYLE.DOT,
                'dash_dot': MSO_LINE_DASH_STYLE.DASH_DOT,
            }
            shape.line.dash_style = line_style_map.get(style['line_style'], MSO_LINE_DASH_STYLE.SOLID)

        # Rotation
        if 'rotation' in style:
            shape.rotation = style['rotation']

    def _apply_connector_styling(self, connector, element: Dict[str, Any]) -> None:
        """Apply styling to connectors (arrows, line styles, etc.).

        For patent diagrams, default is:
        - Black line color
        - 1pt line width
        - Arrow at end
        """
        style = element.get('style', {})

        # Arrow style
        if style.get('arrow_start', False):
            connector.line.begin_arrow_type = 1  # Arrow at start

        if style.get('arrow_end', True):  # Default: arrow at end
            connector.line.end_arrow_type = 1  # Arrow at end

        # Line color - default to BLACK
        line_color = style.get('line_color', '000000')
        rgb = self._hex_to_rgb(line_color)
        connector.line.color.rgb = RGBColor(*rgb)

        # Line width - default to 1pt
        line_width = style.get('line_width', 1.0)
        connector.line.width = Pt(line_width)

        # Line style - default to solid
        if 'line_style' in style:
            line_style_map = {
                'solid': MSO_LINE_DASH_STYLE.SOLID,
                'dash': MSO_LINE_DASH_STYLE.DASH,
                'dot': MSO_LINE_DASH_STYLE.DOT,
                'dash_dot': MSO_LINE_DASH_STYLE.DASH_DOT,
            }
            connector.line.dash_style = line_style_map.get(style['line_style'], MSO_LINE_DASH_STYLE.SOLID)

    def _add_text_to_shape(self, shape, element: Dict[str, Any]) -> None:
        """Add and format text within a shape.

        For patent diagrams, default is:
        - Black text color
        - Center aligned
        """
        text = element.get('text', '')
        text_frame = shape.text_frame
        text_frame.clear()  # Remove default paragraph

        paragraph = text_frame.paragraphs[0]
        paragraph.text = text

        # Text formatting
        text_config = element.get('text_format', {})

        if 'font_size' in text_config:
            paragraph.font.size = Pt(text_config['font_size'])

        if 'font_name' in text_config:
            paragraph.font.name = text_config['font_name']

        if 'bold' in text_config:
            paragraph.font.bold = text_config['bold']

        if 'italic' in text_config:
            paragraph.font.italic = text_config['italic']

        # Text color - default to BLACK
        text_color = text_config.get('color', '000000')
        rgb = self._hex_to_rgb(text_color)
        paragraph.font.color.rgb = RGBColor(*rgb)

        # Text alignment - default to CENTER
        align_map = {
            'left': PP_ALIGN.LEFT,
            'center': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT,
            'justify': PP_ALIGN.JUSTIFY,
        }
        text_align = text_config.get('align', 'center')
        paragraph.alignment = align_map.get(text_align, PP_ALIGN.CENTER)

        # Vertical alignment - default to MIDDLE
        valign_map = {
            'top': MSO_ANCHOR.TOP,
            'middle': MSO_ANCHOR.MIDDLE,
            'bottom': MSO_ANCHOR.BOTTOM,
        }
        vertical_align = text_config.get('vertical_align', 'middle')
        text_frame.vertical_anchor = valign_map.get(vertical_align, MSO_ANCHOR.MIDDLE)

    def _apply_text_formatting(self, text_frame, element: Dict[str, Any]) -> None:
        """Apply text formatting to a text frame."""
        text_config = element.get('text_format', {})
        paragraph = text_frame.paragraphs[0]

        if 'font_size' in text_config:
            paragraph.font.size = Pt(text_config['font_size'])

        if 'font_name' in text_config:
            paragraph.font.name = text_config['font_name']

        if 'bold' in text_config:
            paragraph.font.bold = text_config['bold']

        if 'italic' in text_config:
            paragraph.font.italic = text_config['italic']

        if 'color' in text_config:
            rgb = self._hex_to_rgb(text_config['color'])
            paragraph.font.color.rgb = RGBColor(*rgb)

        if 'align' in text_config:
            align_map = {
                'left': PP_ALIGN.LEFT,
                'center': PP_ALIGN.CENTER,
                'right': PP_ALIGN.RIGHT,
            }
            paragraph.alignment = align_map.get(text_config['align'], PP_ALIGN.LEFT)

    def _get_connection_point(self, shape, side: str) -> Tuple[int, int]:
        """
        Calculate the connection point on a shape's edge.

        Args:
            shape: The shape object
            side: 'top', 'bottom', 'left', 'right', or 'center'

        Returns:
            Tuple of (x, y) coordinates
        """
        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height

        points = {
            'top': (left + width // 2, top),
            'bottom': (left + width // 2, top + height),
            'left': (left, top + height // 2),
            'right': (left + width, top + height // 2),
            'center': (left + width // 2, top + height // 2),
        }

        return points.get(side, points['center'])

    @staticmethod
    def _hex_to_rgb(hex_color: str) -> Tuple[int, int, int]:
        """
        Convert hex color to RGB tuple.

        Args:
            hex_color: Hex string like '#FF0000' or 'FF0000'

        Returns:
            RGB tuple (r, g, b)
        """
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

    def save(self, filepath: str) -> None:
        """
        Save the presentation to a file.

        Args:
            filepath: Path to save the PPTX file
        """
        self.prs.save(filepath)
        logger.info(f"Presentation saved to: {filepath}")


# Convenience function for single-slide generation
def generate_diagram(spec: Dict[str, Any], output_path: str) -> str:
    """
    Generate a diagram from JSON spec and save to file.

    Args:
        spec: Diagram specification
        output_path: Where to save the PPTX file

    Returns:
        Path to saved file
    """
    generator = PPTXDiagramGenerator()
    generator.create_from_json(spec)
    generator.save(output_path)
    return output_path
